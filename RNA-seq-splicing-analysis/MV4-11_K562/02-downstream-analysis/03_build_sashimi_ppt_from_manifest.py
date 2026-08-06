#!/usr/bin/env python3
"""Build a sashimi summary PowerPoint from a manifest or plot-status table.

The preferred manifest contract contains:
- plot_id, cohort, AS_type, geneSymbol, event_display
- plot_output_dir
- ppt_title, ppt_subtitle
- delta_psi_group2_minus_group1

No pandas is required.
"""

import argparse
import csv
import shutil
import subprocess
import tempfile
from collections import Counter, defaultdict
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

REQUIRED_COLUMNS = [
    "plot_id",
    "AS_type",
    "event_display",
    "cohort",
    "geneSymbol",
]

SUCCESS_STATUSES = {
    "completed",
    "completed_multiple_pdfs",
    "skipped_existing",
    "pending",  # permits direct use of the original manifest after PDFs exist
}


def require_cmd(cmd):
    if shutil.which(cmd) is None:
        raise SystemExit(f"Missing required command: {cmd}")


def read_manifest(path):
    with open(path, newline="") as f:
        reader = csv.DictReader(f, delimiter="\t")
        missing = [c for c in REQUIRED_COLUMNS if c not in (reader.fieldnames or [])]
        if missing:
            raise SystemExit("Missing columns in manifest: " + ", ".join(missing))
        return list(reader)


def format_num(value):
    try:
        text = str(value).strip()
        if text == "" or text.upper() in {"NA", "NAN"}:
            return "NA"
        return f"{float(text):.4g}"
    except (TypeError, ValueError):
        return str(value)


def resolve_path(value, manifest_dir):
    path = Path(value)
    if not path.is_absolute():
        path = manifest_dir / path
    return path


def find_pdfs(row, manifest_dir, legacy_output_root=None):
    if row.get("pdf_files"):
        pdfs = [Path(x) for x in row["pdf_files"].split(";") if x]
        pdfs = [x for x in pdfs if x.exists()]
        if pdfs:
            return sorted(pdfs)

    for column in ("plot_output_dir", "expected_pdf_dir"):
        if row.get(column):
            path = resolve_path(row[column], manifest_dir)
            if path.exists():
                pdfs = sorted(path.rglob("*.pdf"))
                if pdfs:
                    return pdfs

    if legacy_output_root:
        root = Path(legacy_output_root)
        candidates = [
            root / row["cohort"] / row["AS_type"] / row["plot_id"],
            root / row["AS_type"] / row["plot_id"],
        ]
        for path in candidates:
            if path.exists():
                pdfs = sorted(path.rglob("*.pdf"))
                if pdfs:
                    return pdfs
    return []


def render_pdf(pdf, out_prefix, dpi=200):
    cmd = ["pdftoppm", "-png", "-r", str(dpi), str(pdf), str(out_prefix)]
    subprocess.run(cmd, check=True)
    return sorted(out_prefix.parent.glob(out_prefix.name + "-*.png"))


def add_textbox(slide, left, top, width, height, text, size=20, bold=False,
                align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    return box


def add_image(slide, img, left, top, width, height):
    with Image.open(img) as im:
        w, h = im.size
    aspect_img = w / h
    aspect_box = width / height
    if aspect_img > aspect_box:
        new_w = width
        new_h = int(width / aspect_img)
        new_left = left
        new_top = top + (height - new_h) // 2
    else:
        new_h = height
        new_w = int(height * aspect_img)
        new_top = top
        new_left = left + (width - new_w) // 2
    slide.shapes.add_picture(str(img), new_left, new_top, width=new_w, height=new_h)


def add_title_slide(prs, blank, run_name, rows):
    slide = prs.slides.add_slide(blank)
    add_textbox(
        slide, Inches(0.7), Inches(1.8), prs.slide_width - Inches(1.4), Inches(1.0),
        f"{run_name}\nSashimi plot summary", size=30, bold=True,
        align=PP_ALIGN.CENTER,
    )
    genes = len({r["geneSymbol"] for r in rows})
    add_textbox(
        slide, Inches(1.0), Inches(3.3), prs.slide_width - Inches(2.0), Inches(1.0),
        f"{genes} genes | {len(rows)} manifest events\n"
        "Displayed ΔPSI convention: group 2 minus group 1",
        size=16, align=PP_ALIGN.CENTER,
    )


def add_summary_slide(prs, blank, rows, plotted_rows, missing_rows):
    slide = prs.slides.add_slide(blank)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.5),
                "Run summary", size=24, bold=True)

    statuses = Counter(r.get("workflow_status", "not_recorded") for r in rows)
    by_type = Counter(r["AS_type"] for r in plotted_rows)
    by_cohort = Counter(r["cohort"] for r in plotted_rows)

    lines = [
        f"Manifest events: {len(rows)}",
        f"Events included in deck: {len(plotted_rows)}",
        f"Events without a usable PDF: {len(missing_rows)}",
        "",
        "Workflow status:",
    ]
    lines.extend(f"  {key}: {value}" for key, value in sorted(statuses.items()))
    lines.extend(["", "Included events by AS type:"])
    lines.extend(f"  {key}: {value}" for key, value in sorted(by_type.items()))
    lines.extend(["", "Included events by cohort:"])
    lines.extend(f"  {key}: {value}" for key, value in sorted(by_cohort.items()))

    add_textbox(slide, Inches(0.7), Inches(1.0), Inches(11.9), Inches(5.9),
                "\n".join(lines), size=16)


def make_subtitle(row):
    if row.get("ppt_subtitle"):
        return row["ppt_subtitle"]

    dpsi = row.get("delta_psi_group2_minus_group1", "")
    if not dpsi:
        # Legacy fallback is clearly labeled as the raw rMATS convention.
        raw = row.get("IncLevelDifference_raw", row.get("IncLevelDifference", ""))
        dpsi_text = f"raw rMATS dPSI(group1-group2)={format_num(raw)}"
    else:
        g1 = row.get("group1_label", "group1")
        g2 = row.get("group2_label", "group2")
        dpsi_text = f"dPSI {g2}-{g1}={format_num(dpsi)}"

    return (
        f"{row['cohort']} | {row['AS_type']} | "
        f"{row.get('feature_label', '')} | "
        f"P={format_num(row.get('PValue', ''))} | "
        f"FDR={format_num(row.get('FDR', ''))} | {dpsi_text}"
    )


def make_ppt(manifest, out_file, dpi=200, skip_missing=False,
             legacy_output_root=None, run_name=None):
    manifest = Path(manifest)
    rows = read_manifest(manifest)
    manifest_dir = manifest.parent

    prs = Presentation()
    prs.slide_width = Inches(13.3)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    resolved = []
    missing = []
    for row in rows:
        status = row.get("workflow_status", "pending")
        if status not in SUCCESS_STATUSES:
            missing.append(row)
            continue
        pdfs = find_pdfs(row, manifest_dir, legacy_output_root)
        if not pdfs:
            missing.append(row)
        else:
            resolved.append((row, pdfs))

    if missing and not skip_missing:
        ids = ", ".join(r["plot_id"] for r in missing[:10])
        raise SystemExit(f"Missing/failed PDFs for {len(missing)} events; examples: {ids}")

    display_run_name = run_name or manifest.parent.parent.name
    add_title_slide(prs, blank, display_run_name, rows)
    add_summary_slide(prs, blank, rows, [x[0] for x in resolved], missing)

    with tempfile.TemporaryDirectory() as tmp_name:
        tmp = Path(tmp_name)
        for row, pdfs in resolved:
            title = row.get("ppt_title") or row["event_display"]
            subtitle = make_subtitle(row)

            for pdf_index, pdf in enumerate(pdfs, 1):
                prefix = tmp / f"{row['plot_id']}_{pdf_index}"
                pngs = render_pdf(pdf, prefix, dpi=dpi)
                for page_index, png in enumerate(pngs, 1):
                    slide = prs.slides.add_slide(blank)
                    add_textbox(
                        slide, Inches(0.3), Inches(0.15),
                        prs.slide_width - Inches(0.6), Inches(0.45),
                        title, size=22, bold=True,
                    )
                    add_textbox(
                        slide, Inches(0.3), Inches(0.62),
                        prs.slide_width - Inches(0.6), Inches(0.38),
                        subtitle, size=10,
                    )
                    add_image(
                        slide, png, Inches(0.3), Inches(1.05),
                        prs.slide_width - Inches(0.6),
                        prs.slide_height - Inches(1.35),
                    )

    out_file = Path(out_file)
    out_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_file)
    print("Saved PPT:", out_file)
    print("Slides:", len(prs.slides))
    print("Missing/failed events omitted:", len(missing))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--manifest", required=True)
    parser.add_argument("--out", required=True)
    parser.add_argument("--dpi", type=int, default=200)
    parser.add_argument("--skip-missing", action="store_true")
    parser.add_argument("--output-root", help="Legacy fallback only")
    parser.add_argument("--run-name")
    args = parser.parse_args()

    require_cmd("pdftoppm")
    make_ppt(
        manifest=args.manifest,
        out_file=args.out,
        dpi=args.dpi,
        skip_missing=args.skip_missing,
        legacy_output_root=args.output_root,
        run_name=args.run_name,
    )


if __name__ == "__main__":
    main()
