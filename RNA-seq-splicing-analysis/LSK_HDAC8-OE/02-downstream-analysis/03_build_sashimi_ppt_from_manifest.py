#!/usr/bin/env python3
"""
Build a PowerPoint from a sashimi manifest (NO pandas required)

Dependencies:
- python-pptx
- pillow
- pdftoppm (poppler)

Usage:
python3 build_sashimi_ppt_from_manifest.py \
  --manifest sashimi_manifest.tsv \
  --output-root output \
  --out sashimi.pptx
"""

import argparse
import csv
import shutil
import subprocess
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


REQUIRED_COLUMNS = [
    "plot_id",
    "AS_type",
    "event_display",
    "cohort",
    "feature_label",
    "FDR",
    "IncLevelDifference",
]


# ---------------------------
# Utils
# ---------------------------
def require_cmd(cmd):
    if shutil.which(cmd) is None:
        raise SystemExit(f"Missing required command: {cmd}")


def ensure_columns(fieldnames):
    missing = [c for c in REQUIRED_COLUMNS if c not in fieldnames]
    if missing:
        raise SystemExit("Missing columns in manifest: " + ", ".join(missing))


def read_manifest(path):
    with open(path, newline="") as f:
        reader = csv.DictReader(f, delimiter="\t")
        ensure_columns(reader.fieldnames)
        return list(reader)


def format_num(x):
    try:
        s = str(x).strip()
        if s == "" or s.upper() in {"NA", "NAN"}:
            return "NA"
        return f"{float(s):.4g}"
    except:
        return str(x)


# ---------------------------
# PDF handling
# ---------------------------
def find_pdfs(output_root, row):
    path = Path(output_root) / row["cohort"] / row["AS_type"] / row["plot_id"]
    if not path.exists():
        return []
    return sorted(path.rglob("*.pdf"))


def render_pdf(pdf, out_prefix, dpi=200):
    cmd = [
        "pdftoppm",
        "-png",
        "-r", str(dpi),
        str(pdf),
        str(out_prefix)
    ]
    subprocess.run(cmd, check=True)
    return sorted(out_prefix.parent.glob(out_prefix.name + "-*.png"))


# ---------------------------
# PPT helpers
# ---------------------------
def add_textbox(slide, left, top, width, height, text, size=20, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold


def add_image(slide, img, left, top, width, height):
    with Image.open(img) as im:
        w, h = im.size

    aspect_img = w / h
    aspect_box = width / height

    if aspect_img > aspect_box:
        new_w = width
        new_h = int(width / aspect_img)
        new_left = left
        new_top = top + (height - new_h) // 2
    else:
        new_h = height
        new_w = int(height * aspect_img)
        new_top = top
        new_left = left + (width - new_w) // 2

    slide.shapes.add_picture(str(img), new_left, new_top, width=new_w, height=new_h)


# ---------------------------
# Main
# ---------------------------
def make_ppt(manifest, output_root, out_file, dpi=200, skip_missing=False):
    rows = read_manifest(manifest)

    prs = Presentation()
    prs.slide_width = Inches(13.3)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory() as tmp:
        tmp = Path(tmp)

        for row in rows:
            pdfs = find_pdfs(output_root, row)

            if not pdfs:
                if skip_missing:
                    print(f"Skip missing: {row['plot_id']}")
                    continue
                else:
                    raise SystemExit(f"No PDF for {row['plot_id']}")

            title = row["event_display"]
            subtitle = (
                f"{row['cohort']} | {row['AS_type']} | "
                f"{row['feature_label']} | "
                f"FDR={format_num(row['FDR'])} | "
                f"dPSI={format_num(row['IncLevelDifference'])}"
            )

            for pdf in pdfs:
                pngs = render_pdf(pdf, tmp / row["plot_id"], dpi=dpi)

                for png in pngs:
                    slide = prs.slides.add_slide(blank)

                    add_textbox(slide, Inches(0.3), Inches(0.2),
                                prs.slide_width - Inches(0.6), Inches(0.4),
                                title, size=22, bold=True)

                    add_textbox(slide, Inches(0.3), Inches(0.6),
                                prs.slide_width - Inches(0.6), Inches(0.3),
                                subtitle, size=10)

                    add_image(slide, png,
                              Inches(0.3), Inches(1.0),
                              prs.slide_width - Inches(0.6),
                              prs.slide_height - Inches(1.3))

    prs.save(out_file)
    print("Saved PPT:", out_file)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--manifest", required=True)
    parser.add_argument("--output-root", required=True)
    parser.add_argument("--out", required=True)
    parser.add_argument("--dpi", type=int, default=200)
    parser.add_argument("--skip-missing", action="store_true")

    args = parser.parse_args()

    require_cmd("pdftoppm")

    make_ppt(
        manifest=args.manifest,
        output_root=args.output_root,
        out_file=args.out,
        dpi=args.dpi,
        skip_missing=args.skip_missing
    )


if __name__ == "__main__":
    main()